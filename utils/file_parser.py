import fitz
from pptx import Presentation
from docx import Document

def parse_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def parse_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text +=shape.text + "\n"
    return text

def parse_doc(path):
    doc = Document(path)
    text =""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text